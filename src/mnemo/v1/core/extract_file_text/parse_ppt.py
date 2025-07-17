import re

from pptx import Presentation

from mnemo.v1.core.extract_file_text.parse_doc_utils import (
    check_line_validity,
    clean_whitespace,
)


def extract_all_texts_in_ppt(file_path: str) -> str:
    """
    parse a PPTX file and extract text

    Args:
        file_path (str): path to the PPTX file
    Returns:
        str: extracted text from the PPTX file, cleaned and formatted
    """

    prs = Presentation(file_path)
    paras = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and check_line_validity(text):
                    text = clean_whitespace(text)
                else:
                    continue
                paras.append(text)
    text = "\n".join(paras)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def parse_ppt_by_slide(file_path: str) -> str:
    """
    parse a PPTX file and extract text by slide

    Args:
        file_path (str): path to the PPTX file
    Returns:
        list[str]: list of strings, each string is the text of a slide
    """
    prs = Presentation(file_path)
    all_slides = []

    for slide in prs.slides:
        silde_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and check_line_validity(text):
                    text = clean_whitespace(text)
                else:
                    continue
                silde_text.append(text)
        if silde_text:
            slide_text = "\n".join(silde_text)
            all_slides.append(slide_text)
    return all_slides


def parse_ppt_to_get_headers(ppt_path: str, min_header_fontsize: int = 14, min_header_rank: int = 3) -> list[str]:
    """
    parse a PPTX file and extract headers

    Args:
        ppt_path (str): path to the PPTX file
        min_header_fontsize (int): minimum font size to consider a text as header
        min_header_rank (int): minimum rank of the header to consider it as header
    Returns:
        list[str]: list of headers extracted from the PPTX file
    """

    prs = Presentation(ppt_path)
    all_headers = []

    for slide in prs.slides:
        slide_headers = []
        shape_texts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue  # skip shapes without text frames

            for para in shape.text_frame.paragraphs:
                if not para.text or not check_line_validity(para.text):
                    continue

                for run in para.runs:  # get the font size from the first run
                    if run.font.size and run.font.size.pt:
                        font_size = run.font.size.pt
                        break
                    else:
                        font_size = 10  # default font size if not specified

                shape_texts.append({'text': clean_whitespace(para.text), 'top': shape.top, 'font_size': font_size, })

        shape_texts.sort(key=lambda x: x['top'])

        for i, item in enumerate(shape_texts):
            text = item['text']
            is_header = False

            if i <= min_header_rank:
                if item['font_size'] >= min_header_fontsize:
                    is_header = True
            elif re.match(r'^\d+[\.\s]', text):
                is_header = True
            elif text.startswith(('• ', '- ', '◦ ', '▪ ')):
                is_header = True

            if is_header:
                slide_headers.append(text)

        if slide_headers:
            slide_header_text = "\n".join(slide_headers)
            all_headers.append(slide_header_text)

    return all_headers
